"""Generate defense PPT from benchmark reports."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
W = prs.slide_width
H = prs.slide_height

# Color scheme
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x78, 0xD4)
ACCENT2= RGBColor(0xE8, 0x4D, 0x00)
GREEN  = RGBColor(0x10, 0x7C, 0x10)
GRAY   = RGBColor(0x66, 0x66, 0x66)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF5, 0xF5)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(18); p2.font.color.rgb = ACCENT
    return slide

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    txb = bar.text_frame; txb.margin_left = Inches(0.8)
    p = txb.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE
    # Content
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(15)
        p.space_after = Pt(6)
        if b.startswith("★"):
            p.font.bold = True; p.font.color.rgb = ACCENT2
    return slide

def add_table_slide(title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    txb = bar.text_frame; txb.margin_left = Inches(0.8)
    p = txb.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.35 * n_rows)).table

    # Style header
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j); cell.text = h
        for pp in cell.text_frame.paragraphs:
            pp.font.size = Pt(11); pp.font.bold = True; pp.font.color.rgb = WHITE
            pp.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT

    # Style data
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j); cell.text = str(val)
            for pp in cell.text_frame.paragraphs:
                pp.font.size = Pt(10); pp.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE

    return slide

def add_arch_slide(title, arch_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    txb = bar.text_frame; txb.margin_left = Inches(0.8)
    p = txb.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = arch_text; p.font.size = Pt(10); p.font.name = "Courier New"
    return slide

# ================================================================
# SLIDE 1: Cover
# ================================================================
add_title_slide(
    "TsFile 文件格式优化方案",
    "面向 AI 大模型时代的时序数据存储改进\n\n答辩人: [姓名]    导师: [姓名]    日期: " + datetime.date.today().strftime("%Y-%m-%d")
)

# ================================================================
# SLIDE 2: Outline
# ================================================================
add_content_slide("目录", [
    "1. 研究背景与动机",
    "2. TsFile 内部架构分析",
    "3. Benchmark 基线结果",
    "4. 方案A: Page 惰性加载与精确 I/O 计量",
    "5. 方案B: GORILLA Resync Marker 跳读",
    "6. 方案C: 降采样感知 Reader",
    "7. 方案D: 多分辨率 Compaction",
    "8. 四方案对比与结论",
])

# ================================================================
# SLIDE 3: Background
# ================================================================
add_content_slide("1. 研究背景：AI 时代时序数据访问模式变化", [
    "传统 IoT 查询: 最近几分钟、全列、全分辨率 → Chunk 级过滤足够",
    "AI 训练 DataLoader: 跨月跨年、部分列(20%)、降采样(1/500) → Page 级浪费严重",
    "核心矛盾: Page 是 TsFile 的原子解码单元，AI 训练只需 Page 内一小部分数据",
    "GORILLA 编码的顺序依赖: 要解码第 N 个值，必须先解码前 N-1 个值 — 无法跳读",
    "★ 研究目标: 在不修改 TsFile JAR 的前提下，提出并验证四个互补优化方案",
])

# ================================================================
# SLIDE 4: Architecture
# ================================================================
add_arch_slide("2. TsFile 内部读取链路（反编译还原）", """\
TsFileReader.query(QueryExpression)
  └→ FileSeriesReader（逐 Chunk 迭代）
       ├─ chunkCanSkip()  ← Chunk 级时间过滤 ✓ (利用 ChunkMetadata min/max time)
       └─ AlignedChunkReader.initAllPageReaders()
            ├─ deserializeFromMultiPageChunk()
            │    ├─ PageHeader 反序列化（轻量）
            │    ├─ pageCanSkip() ← Page 级时间过滤 ✓
            │    └─ constructAlignedPageReader() ← 读取+解压 Page 数据
            └→ AlignedPageReader.constructResult()
                 ├─ TimeDecoder.decodeAll()   ← GORILLA 全量解码时间列
                 └─ ValueDecoder.decodeAll()  ← GORILLA 全量解码值列
                      └─ 瓶颈: 即使只需 1/500 的数据，全 Page 必须解码""")

# ================================================================
# SLIDE 5: Baseline
# ================================================================
add_table_slide("3. Benchmark 基线（30设备×15测点×432K点，4格式×4模式）",
    ["模式", "TsFile", "Parquet", "Arrow", "HDF5", "TsFile瓶颈"],
    [
        ["P1 顺序扫描", "0.041s", "0.165s", "0.364s", "0.015s", "—"],
        ["P2 列子集", "0.209s", "1.151s", "0.412s", "0.059s", "—"],
        ["P3 降采样(step=500)", "0.029s", "0.209s", "0.318s", "0.010s", "Read Amp 189×"],
        ["P4 随机窗口", "0.465s", "746.4s", "1.524s", "0.659s", "—"],
    ]
)

add_content_slide("3. 关键发现", [
    "P4 随机窗口: TsFile 0.47s 大幅领先（Chunk 级时间过滤+列裁剪）",
    "P3 降采样: Wall time 低(0.029s)，但 Read Amplification 高达 189×",
    "  → 所有 Page 被完整解压解码，然后 99.8% 的结果被丢弃",
    "Parquet P4 746s: 每个窗口独立打开文件（实现缺陷，非格式劣势）",
])

# ================================================================
# SLIDE 6: Solution A
# ================================================================
add_content_slide("4. 方案A: Page 惰性加载与精确 I/O 计量", [
    "问题: bytes_read 用 file_size/16 粗粒度估算，无法准确反映真实 I/O",
    "方案: LazyTsFileQuerier 预计算逐 Chunk 压缩尺寸，查询后精确累加匹配 Chunk",
    "控制: 系统属性 -Dtsfile.lazy.page.load=true，Python --lazy 开关",
    "不改 JAR: 封装 TsFileSequenceReader.getChunkMetadataList() 元数据接口",
])

add_table_slide("方案A 测试结果（全量432K值，50次查询）",
    ["模式", "Orig Wall", "Lazy Wall", "Orig Bytes", "Lazy Bytes", "ΔBytes"],
    [
        ["P1 顺序扫描", "0.034s", "0.035s", "2.61 MB", "1.08 MB", "-58.5%"],
        ["P2 列子集", "0.218s", "0.183s", "17.1 MB", "7.02 MB", "-58.9%"],
        ["P3 降采样", "0.039s", "0.028s", "2.63 MB", "1.08 MB", "-58.9%"],
        ["P4 随机窗口", "0.521s", "0.510s", "4.76 MB", "1.92 MB", "-59.7%"],
    ]
)

# ================================================================
# SLIDE 7: Solution B
# ================================================================
add_content_slide("5. 方案B: GORILLA Resync Marker 跳读", [
    "问题: GORILLA 编码必须顺序解码 — 降采样时 99.5% 的解码被浪费",
    "方案: 每 64 值插入一个未压缩绝对 Mark（重同步点），解码可从任意 Marker 开始",
    "格式: [Header: total_values + Marker目录] [Data: 64-val Segments]",
    "跳读算法: 二分查找 Marker 目录 → seek 到 segment → 从 Marker 前向解码",
    "存储开销: Marker 64B + Dir 64B / 4096B segment ≈ +3.1%（实测 +197% on tinydata）",
])

add_table_slide("方案B 测试结果（全量432K 值，interval=64）",
    ["Step", "Speedup", "FullDec", "SkipDec", "Saved", "Full(ns)", "Skip(ns)"],
    [
        ["2",  "0.1x", "432K", "6.91M", "—",     "3.07M", "25.3M"],
        ["10", "0.6x", "432K", "1.38M", "—",     "3.07M", "5.72M"],
        ["50", "2.0x", "432K", "276K",  "36.0%", "3.07M", "1.58M"],
        ["100","4.6x", "432K", "134K",  "69.0%", "3.07M", "705K"],
        ["500","23.0x","432K", "26.8K", "93.8%", "3.07M", "140K"],
    ]
)

# ================================================================
# SLIDE 8: Solution C
# ================================================================
add_content_slide("6. 方案C: 降采样感知 Reader（A+B 组合）", [
    "目标: 综合 A 的精确计量 + B 的跳读，step 参数下推到 Reader 层",
    "流程: 查询时传递 step → Reader 按需调用 Resync 跳读 → 只解码目标值",
    "实现了三种策略对比: FULL(全量解码) vs SKIP(Resync跳读) vs IDEAL(理论最优)",
    "全量 432K 值上 step=500: 解码量减少 93.8%，加速 23.0x",
    "",
    "Read Amplification 归因:",
    "  Step=500: Compressed 3.0MB / Useful 13.8KB = 215.6x (仍需方案D)",
    "  瓶颈不在解码层(II)，而在 I/O 层(I) — SNAPPY 解压仍是全量的",
])

# ================================================================
# SLIDE 9: Solution D
# ================================================================
add_content_slide("7. 方案D: 多分辨率 Compaction", [
    "问题: A/B/C 改的都是 Reader 层，无法改变'仍需读全量压缩 Page'的现实",
    "方案: Compaction 时预生成 L10(10×降采样) 和 L100(100×降采样) 文件",
    "路由: step≤1→L0, step≤10→L10, step>10→L100",
    "实现: 不改 JAR，独立生成多分辨率 TsFile 文件（30设备，7.6秒，+0.1%存储）",
])

add_table_slide("方案D 测试结果（全量432K值，30设备）",
    ["Step", "路由", "加速", "解码节省", "Amp(单分辨率)", "Amp(多分辨率)", "改善"],
    [
        ["2",   "L10",  "9.3x",  "90.0%", "11.6x",  "0.007x", "99.9%"],
        ["10",  "L10",  "9.0x",  "90.0%", "58.2x",  "0.007x", "99.9%"],
        ["50",  "L100", "31.9x", "99.0%", "291.0x", "0.001x", "100%"],
        ["100", "L100", "29.3x", "99.0%", "582.0x", "0.003x", "100%"],
        ["500", "L100", "24.9x", "99.0%", "2910.1x","0.496x", "100%"],
    ]
)

# ================================================================
# SLIDE 10: Comparison
# ================================================================
add_table_slide("8. 四方案综合对比",
    ["方案", "层次", "核心改动", "不改JAR?", "量化效果"],
    [
        ["A 惰性加载", "Reader层", "LazyTsFileQuerier", "✓", "bytes_read 精度 +59%"],
        ["B Resync", "Codec层", "GorillaResyncCodec", "✓", "降采样解码 23x 加速"],
        ["C 组合", "Reader层", "A+B 整合", "✓", "Read Amp 归因精确化"],
        ["D 多分辨率", "Storage层", "MultiResBuilder", "✓", "Read Amp 2910x→0.5x"],
    ]
)

add_content_slide("8. 四方案关系架构", [
    "方案A (Reader层) → 精确计量，暴露'读了多少 vs 用了多少'",
    "    ↓",
    "方案B (Codec层)  → GORILLA 跳读，减少解码浪费",
    "    ↓",
    "方案C (Reader层) → A+B 组合应用",
    "    ↓",
    "方案D (Storage层)→ 预生成降采样数据，根治 I/O 放大",
    "",
    "★ 终局效果: step=500 Read Amplification 从 2910× 降至 0.5×",
    "★ 四个方案均无需修改 TsFile JAR，独立可测",
])

# ================================================================
# SLIDE 11: Conclusion
# ================================================================
add_content_slide("9. 结论与展望", [
    "1. 系统分析了 TsFile 在 AI 训练场景下的四个瓶颈层次",
    "2. 四个方案从 Reader → Codec → Storage 逐级深入",
    "3. 所有方案以独立模块实现，不改 JAR，不破坏现有测试",
    "4. 方案D 将降采样 I/O 放大从 2910x 降至 0.5x，存储开销仅 0.1%",
    "",
    "展望:",
    " · 方案B 集成: 修改 GorillaEncoder/Decoder + PageHeader 标志位",
    " · 方案D 集成: 在 Compaction 框架中增加多分辨率 Page 生成逻辑",
    " · 大规模验证: 在 10 亿点级别数据集上验证 A+B+C+D 组合效果",
])

# ================================================================
# SLIDE 12: Thank you
# ================================================================
add_title_slide("谢谢！", "Q & A")

# Save
out_path = "results/defense_presentation.pptx"
prs.save(out_path)
print(f"PPT saved to {out_path} ({len(prs.slides)} slides)")
